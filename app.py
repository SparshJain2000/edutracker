import sys
import os
# from flask import Flask, render_template
from nltk.stem import WordNetLemmatizer
from nltk.tokenize import RegexpTokenizer
from pptx import Presentation
import numpy as np
# from sklearn.feature_extraction.text import TfidfVectorizer
from nltk.tokenize import word_tokenize, sent_tokenize
from nltk.corpus import stopwords
stop_words = stopwords.words('english')

# your file location
# filename = r"./vec.pptx"
prs = Presentation(os.path.join("./",  sys.argv[1]))
slide_titles = []
for slide in prs.slides:
    if slide.shapes.title is None:
        continue
    titles = slide.shapes.title.text
    slide_titles.append(titles)
# print(slide_titles)
slide_titles = np.array(slide_titles)
# print(slide_titles)
#from nltk.tokenize import RegexpTokenizer
e = []
tokenizer = RegexpTokenizer(r'\w+')
for i in slide_titles:
    d = tokenizer.tokenize(i)
    content = (' '.join([j for j in d]))
    e.append(content)
main_data = []
for i in e:
    w = i.split()
    content = (' '.join([j for j in w if not (j.isdigit())]))
    if i == '':
        continue
    main_data.append(content)
data = []
for i in main_data:
    w = i.split()
    content = (' '.join([j for j in w if j.lower() not in stop_words]))
    data.append(content)
data = np.array(data)
data = np.char.lower(data)
# print(data)
#from nltk.stem import WordNetLemmatizer
lemmatizer = WordNetLemmatizer()
final_data = []
for i in data:
    w = i.split()
    words_lemma = [lemmatizer.lemmatize(word) for word in w]
    content = (' '.join([j for j in words_lemma]))
    final_data.append(content)
final_data = np.unique(final_data)
w = ""
for i in final_data:
    w = w+","+i
print(w)
